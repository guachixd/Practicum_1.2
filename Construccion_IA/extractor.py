"""
extractor.py — Extractor de Texto

Corresponde al contenedor "Extractor de Texto" de tu diagrama C4.
Lee PDF, Word (.docx) y PowerPoint (.pptx).

Cambio importante: ya NO busca archivos por nombre (plan_docente, guia_estudio, etc.).
En vez de eso, lee TODOS los archivos que estén directamente en la raíz de la carpeta
de material (sin importar cómo los haya nombrado el docente) y los combina como
"material general del curso". Esto evita depender de una convención de nombres.

Nota de alcance: los documentos que resultan ser solo imágenes (PDFs escaneados
sin capa de texto) se detectan y se saltan con una advertencia, en vez de
aplicarles OCR — decisión tomada para simplificar el alcance del proyecto.
"""

import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation

# Si un PDF tiene menos caracteres que esto por página en promedio,
# lo consideramos "probablemente escaneado / solo imagen" y lo saltamos.
UMBRAL_MIN_CARACTERES_POR_PAGINA = 30


def extraer_texto_pdf(ruta: str) -> str | None:
    """
    Extrae el texto de un PDF. Devuelve None si el documento parece ser
    solo imágenes (sin capa de texto real), para que el llamador lo salte.
    """
    doc = fitz.open(ruta)
    texto_paginas = []
    for pagina in doc:
        texto_paginas.append(pagina.get_text())
    doc.close()

    texto_completo = "\n".join(texto_paginas)
    promedio_por_pagina = len(texto_completo) / max(len(texto_paginas), 1)

    if promedio_por_pagina < UMBRAL_MIN_CARACTERES_POR_PAGINA:
        print(f"  [!] '{os.path.basename(ruta)}' parece ser solo imágenes (poco texto extraíble). Se omite.")
        return None

    return texto_completo.strip()


def extraer_texto_docx(ruta: str) -> str:
    """Extrae el texto de un documento Word, incluyendo texto de tablas."""
    documento = docx.Document(ruta)
    partes = [p.text for p in documento.paragraphs if p.text.strip()]

    for tabla in documento.tables:
        for fila in tabla.rows:
            texto_fila = " | ".join(celda.text.strip() for celda in fila.cells)
            if texto_fila.strip(" |"):
                partes.append(texto_fila)

    return "\n".join(partes).strip()


def extraer_texto_pptx(ruta: str) -> str:
    """Extrae el texto de todas las diapositivas de un PowerPoint."""
    presentacion = Presentation(ruta)
    partes = []

    for i, diapositiva in enumerate(presentacion.slides, start=1):
        texto_diapositiva = []
        for forma in diapositiva.shapes:
            if forma.has_text_frame:
                for parrafo in forma.text_frame.paragraphs:
                    texto_parrafo = "".join(run.text for run in parrafo.runs)
                    if texto_parrafo.strip():
                        texto_diapositiva.append(texto_parrafo)

        if texto_diapositiva:
            partes.append(f"[Diapositiva {i}]\n" + "\n".join(texto_diapositiva))

    return "\n\n".join(partes).strip()


def extraer_texto(ruta: str) -> str | None:
    """
    Dispatcher: detecta la extensión del archivo y llama al extractor
    correspondiente. Devuelve None si el archivo se omite (ej. PDF escaneado)
    o si la extensión no está soportada.
    """
    extension = os.path.splitext(ruta)[1].lower()

    try:
        if extension == ".pdf":
            return extraer_texto_pdf(ruta)
        elif extension == ".docx":
            return extraer_texto_docx(ruta)
        elif extension == ".pptx":
            return extraer_texto_pptx(ruta)
        else:
            print(f"  [!] Extensión no soportada, se omite: {ruta}")
            return None
    except Exception as e:
        print(f"  [X] Error extrayendo '{os.path.basename(ruta)}': {e}")
        return None


def procesar_material_general(ruta_material: str) -> str:
    """
    Lee TODOS los archivos que estén directamente en la raíz de la carpeta de
    material (sin entrar a las subcarpetas unidad_XX), sin importar cómo se
    llamen, y los combina en un solo texto de contexto general del curso.
    """
    partes = []
    for archivo in sorted(os.listdir(ruta_material)):
        ruta_completa = os.path.join(ruta_material, archivo)
        if not os.path.isfile(ruta_completa):
            continue  # las carpetas unidad_XX se procesan aparte, en procesar_material_curso

        texto = extraer_texto(ruta_completa)
        if texto:
            partes.append(f"--- {archivo} ---\n{texto}")
            print(f"  -> {archivo}: {len(texto)} caracteres extraídos (material general)")

    return "\n\n".join(partes)


def procesar_material_curso(ruta_material: str) -> dict:
    """
    Recorre la carpeta de material del curso y devuelve un diccionario con:
      - material_general: texto combinado de todos los archivos en la raíz
        (plan docente, guías, sílabo, o lo que sea que haya puesto el docente,
        sin importar el nombre del archivo)
      - unidades: { 1: "texto combinado de todos los archivos de unidad_01", ... }

    Estructura de carpetas esperada:
        material/
        ├── (cualquier archivo: plan docente, guía, sílabo, etc. — cualquier nombre)
        ├── unidad_01/
        │   ├── diapositivas.pptx
        │   └── lectura.pdf
        ├── unidad_02/
        │   └── ...
        └── ... hasta unidad_08/
    """
    resultado = {
        "material_general": "",
        "unidades": {n: "" for n in range(1, 9)}
    }

    print("Procesando material general (todos los archivos en la raíz, sin importar el nombre)...")
    resultado["material_general"] = procesar_material_general(ruta_material)

    print("\nProcesando materiales por unidad...")
    for n in range(1, 9):
        carpeta_unidad = os.path.join(ruta_material, f"unidad_{n:02d}")
        if not os.path.isdir(carpeta_unidad):
            print(f"  [!] No existe la carpeta unidad_{n:02d}, se deja vacía")
            continue

        textos_unidad = []
        for archivo in sorted(os.listdir(carpeta_unidad)):
            ruta_completa = os.path.join(carpeta_unidad, archivo)
            if os.path.isfile(ruta_completa):
                texto = extraer_texto(ruta_completa)
                if texto:
                    textos_unidad.append(f"--- {archivo} ---\n{texto}")

        resultado["unidades"][n] = "\n\n".join(textos_unidad)
        print(f"  unidad_{n:02d}: {len(textos_unidad)} documento(s) procesado(s), "
              f"{len(resultado['unidades'][n])} caracteres totales")

    return resultado


# ==========================================
# BLOQUE DE PRUEBA LOCAL
# ==========================================
if __name__ == "__main__":
    ruta_material = os.path.join(os.path.dirname(__file__), "..", "data", "material")
    material = procesar_material_curso(ruta_material)

    print("\n--- RESUMEN DE EXTRACCIÓN ---")
    print(f"Material general: {len(material['material_general'])} caracteres")
    for n, texto in material["unidades"].items():
        print(f"Unidad {n}: {len(texto)} caracteres")
