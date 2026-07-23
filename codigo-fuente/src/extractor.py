"""
extractor.py — Extractor de Texto

Corresponde al contenedor "Extractor de Texto" de tu diagrama C4.
Lee PDF, Word (.docx) y PowerPoint (.pptx).

Cambio de estructura (v2): ya NO se espera una carpeta por unidad
(unidad_01/, unidad_02/, ...). El docente simplemente sube TODOS los
archivos —plan docente, diapositivas, lecturas, guías— sueltos en una
sola carpeta de material, sin ninguna convención de nombres ni de
subcarpetas. La asignación de cada archivo a su unidad correspondiente
ya no se decide aquí por estructura de carpetas: se decide más adelante,
por contenido, en base_conocimiento.py (clasificador contra el plan
docente).

Este módulo se encarga solo de EXTRAER y VALIDAR el texto de cada
archivo individualmente, devolviendo una lista de "documentos" con su
metadata (nombre, tipo, texto limpio, idioma detectado, advertencias),
para que el resto del pipeline decida qué hacer con cada uno.
"""

import os
import re
import unicodedata

import fitz  # PyMuPDF
import docx
from docx.opc.exceptions import PackageNotFoundError
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

try:
    from langdetect import detect, DetectorFactory, LangDetectException
    DetectorFactory.seed = 0  # resultados reproducibles
    _LANGDETECT_DISPONIBLE = True
except ImportError:
    _LANGDETECT_DISPONIBLE = False

EXTENSIONES_SOPORTADAS = {".pdf", ".docx", ".pptx"}

# Si un PDF tiene menos caracteres que esto por página en promedio,
# lo consideramos "probablemente escaneado / solo imagen" y lo saltamos.
UMBRAL_MIN_CARACTERES_POR_PAGINA = 30

# Umbral general: si después de limpiar el texto queda menos que esto,
# se considera "sin contenido útil" (aplica a docx y pptx también, no
# solo a PDF).
UMBRAL_MIN_CARACTERES_UTIL = 40

# Si el idioma detectado del documento no está en esta lista, se marca
# con una advertencia (pero el texto igual se conserva: la asignatura
# puede tener bibliografía en inglés, por ejemplo).
IDIOMAS_ESPERADOS = {"es"}


# ==========================================
# LIMPIEZA / VALIDACIÓN DE TEXTO
# ==========================================

def limpiar_texto(texto: str) -> str:
    """
    Normaliza el texto extraído para que sea seguro pasarlo a Gemini y
    a python-docx: quita caracteres de control invisibles (incluyendo
    bytes nulos que a veces aparecen en documentos con codificación mal
    declarada), normaliza a NFC, colapsa espacios en blanco repetidos y
    recorta espacios al inicio/fin de cada línea.
    """
    if not texto:
        return ""

    texto = unicodedata.normalize("NFC", texto)

    # Quita caracteres de control excepto salto de línea y tab.
    texto = "".join(
        c for c in texto
        if c in ("\n", "\t") or unicodedata.category(c)[0] != "C"
    )

    # Colapsa espacios/tabs repetidos, conserva saltos de línea simples.
    texto = re.sub(r"[ \t]+", " ", texto)
    texto = re.sub(r"\n[ \t]+", "\n", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)

    return texto.strip()


def detectar_idioma(texto: str) -> str | None:
    """
    Detecta el idioma dominante del texto. Devuelve el código ISO
    (ej. "es", "en") o None si no se pudo detectar o la librería no
    está instalada. No es bloqueante: solo informativo para el reporte
    de cobertura y para que el docente sepa si subió algo en otro idioma.
    """
    if not _LANGDETECT_DISPONIBLE or not texto or len(texto) < 20:
        return None
    try:
        return detect(texto)
    except LangDetectException:
        return None


# ==========================================
# EXTRACTORES POR TIPO DE ARCHIVO
# ==========================================

def extraer_texto_pdf(ruta: str) -> tuple[str | None, list[str]]:
    """
    Extrae el texto de un PDF, incluyendo el contenido de tablas cuando
    PyMuPDF logra reconocerlas como tal. Devuelve (texto, advertencias).
    texto es None si el documento parece ser solo imágenes (sin capa de
    texto real) o si está cifrado y no se puede leer.
    """
    advertencias = []
    try:
        doc = fitz.open(ruta)
    except Exception as e:
        return None, [f"No se pudo abrir el PDF (posiblemente corrupto): {e}"]

    if doc.is_encrypted:
        # Intenta abrir sin contraseña (algunos PDFs están "cifrados"
        # solo para restringir edición, no para restringir lectura).
        if not doc.authenticate(""):
            doc.close()
            return None, ["El PDF está protegido con contraseña y no se pudo leer."]
        advertencias.append("El PDF estaba protegido; se pudo leer igualmente sin contraseña.")

    texto_paginas = []
    for pagina in doc:
        texto_paginas.append(pagina.get_text())

        # Tablas: PyMuPDF puede detectar tablas estructuradas en la página.
        try:
            tablas = pagina.find_tables()
            for tabla in tablas.tables:
                filas = tabla.extract()
                texto_tabla = "\n".join(
                    " | ".join(str(celda).strip() if celda else "" for celda in fila)
                    for fila in filas
                )
                if texto_tabla.strip(" |\n"):
                    texto_paginas.append(f"[Tabla]\n{texto_tabla}")
        except Exception:
            pass  # detección de tablas es un plus, no crítico

    doc.close()

    texto_completo = limpiar_texto("\n".join(texto_paginas))
    promedio_por_pagina = len(texto_completo) / max(len(texto_paginas), 1)

    if promedio_por_pagina < UMBRAL_MIN_CARACTERES_POR_PAGINA:
        return None, advertencias + [
            "El PDF parece ser solo imágenes/escaneado (muy poco texto extraíble). Se omite."
        ]

    return texto_completo, advertencias


def extraer_texto_docx(ruta: str) -> tuple[str | None, list[str]]:
    """Extrae el texto de un documento Word, incluyendo texto de tablas."""
    try:
        documento = docx.Document(ruta)
    except PackageNotFoundError:
        return None, ["El archivo .docx está corrupto o no es un Word válido."]
    except Exception as e:
        return None, [f"No se pudo abrir el .docx: {e}"]

    partes = [p.text for p in documento.paragraphs if p.text.strip()]

    for tabla in documento.tables:
        for fila in tabla.rows:
            texto_fila = " | ".join(celda.text.strip() for celda in fila.cells)
            if texto_fila.strip(" |"):
                partes.append(texto_fila)

    return limpiar_texto("\n".join(partes)), []


def extraer_texto_pptx(ruta: str) -> tuple[str | None, list[str]]:
    """
    Extrae el texto de todas las diapositivas de un PowerPoint, incluyendo
    tablas, notas del orador y texto dentro de agrupaciones de formas.
    """
    try:
        presentacion = Presentation(ruta)
    except PptxPackageNotFoundError:
        return None, ["El archivo .pptx está corrupto o no es un PowerPoint válido."]
    except Exception as e:
        return None, [f"No se pudo abrir el .pptx: {e}"]

    partes = []

    def _texto_de_formas(formas):
        fragmentos = []
        for forma in formas:
            if forma.shape_type == 6:  # GROUP: recorre recursivamente
                fragmentos.extend(_texto_de_formas(forma.shapes))
                continue

            if forma.has_text_frame:
                for parrafo in forma.text_frame.paragraphs:
                    texto_parrafo = "".join(run.text for run in parrafo.runs)
                    if texto_parrafo.strip():
                        fragmentos.append(texto_parrafo)

            if forma.has_table:
                for fila in forma.table.rows:
                    texto_fila = " | ".join(celda.text.strip() for celda in fila.cells)
                    if texto_fila.strip(" |"):
                        fragmentos.append(f"[Tabla] {texto_fila}")
        return fragmentos

    for i, diapositiva in enumerate(presentacion.slides, start=1):
        texto_diapositiva = _texto_de_formas(diapositiva.shapes)

        if diapositiva.has_notes_slide:
            texto_notas = diapositiva.notes_slide.notes_text_frame.text
            if texto_notas.strip():
                texto_diapositiva.append(f"[Notas del orador] {texto_notas.strip()}")

        if texto_diapositiva:
            partes.append(f"[Diapositiva {i}]\n" + "\n".join(texto_diapositiva))

    return limpiar_texto("\n\n".join(partes)), []


# ==========================================
# DISPATCHER Y ARMADO DEL RESULTADO
# ==========================================

def extraer_documento(ruta: str) -> dict:
    """
    Extrae y valida un archivo individual. Devuelve un dict "documento"
    con toda la metadata que el resto del pipeline necesita, sin importar
    si el archivo terminó siendo usable o no (queda registrado igual,
    con sus advertencias, para trazabilidad).
    """
    nombre = os.path.basename(ruta)
    extension = os.path.splitext(ruta)[1].lower()

    documento = {
        "nombre": nombre,
        "ruta": ruta,
        "tipo": extension.lstrip("."),
        "texto": "",
        "num_caracteres": 0,
        "idioma": None,
        "utilizable": False,
        "advertencias": [],
    }

    if extension not in EXTENSIONES_SOPORTADAS:
        documento["advertencias"].append(
            f"Extensión no soportada ({extension}); tipos válidos: PDF, DOCX, PPTX. Se omite."
        )
        print(f"  [!] {nombre}: extensión no soportada, se omite.")
        return documento

    try:
        if extension == ".pdf":
            texto, advertencias = extraer_texto_pdf(ruta)
        elif extension == ".docx":
            texto, advertencias = extraer_texto_docx(ruta)
        else:  # .pptx
            texto, advertencias = extraer_texto_pptx(ruta)
    except Exception as e:
        documento["advertencias"].append(f"Error inesperado al procesar el archivo: {e}")
        print(f"  [X] Error extrayendo '{nombre}': {e}")
        return documento

    documento["advertencias"].extend(advertencias)

    if texto is None:
        print(f"  [!] {nombre}: sin texto utilizable. {'; '.join(advertencias) if advertencias else ''}")
        return documento

    if len(texto) < UMBRAL_MIN_CARACTERES_UTIL:
        documento["advertencias"].append(
            "El archivo se pudo leer pero tiene muy poco contenido de texto útil."
        )
        print(f"  [!] {nombre}: contenido insuficiente ({len(texto)} caracteres).")
        documento["texto"] = texto
        documento["num_caracteres"] = len(texto)
        return documento

    idioma = detectar_idioma(texto)
    if idioma and idioma not in IDIOMAS_ESPERADOS:
        documento["advertencias"].append(
            f"Contenido detectado en idioma distinto al español ({idioma}); se incluye igual, "
            f"pero revisa si es intencional (ej. bibliografía en inglés)."
        )

    documento["texto"] = texto
    documento["num_caracteres"] = len(texto)
    documento["idioma"] = idioma
    documento["utilizable"] = True

    etiqueta_idioma = f", idioma: {idioma}" if idioma else ""
    print(f"  -> {nombre}: {len(texto)} caracteres extraídos{etiqueta_idioma}")

    return documento


def procesar_material_curso(ruta_material: str) -> dict:
    """
    Lee TODOS los archivos que estén en la carpeta de material del curso
    (sin ninguna convención de subcarpetas por unidad) y devuelve un
    diccionario con:
      - documentos: lista de dicts, uno por archivo, con su texto y
        metadata (ver extraer_documento). La clasificación de qué
        documento pertenece a qué unidad se hace después, en
        base_conocimiento.py, en base al contenido y al plan docente.
      - advertencias_globales: problemas generales (carpeta vacía, etc.)

    Estructura de carpetas esperada (ya NO hay unidad_01/, unidad_02/...):
        material/
        ├── plan_docente.pdf        (o el nombre que le haya puesto el docente)
        ├── diapositivas_tema1.pptx
        ├── lectura_complementaria.docx
        └── ... cualquier otro archivo pdf/docx/pptx, sin orden ni nombre fijo
    """
    resultado = {"documentos": [], "advertencias_globales": []}

    if not os.path.isdir(ruta_material):
        resultado["advertencias_globales"].append(
            f"La carpeta de material '{ruta_material}' no existe."
        )
        print(f"  [X] La carpeta '{ruta_material}' no existe.")
        return resultado

    archivos = sorted(
        f for f in os.listdir(ruta_material)
        if os.path.isfile(os.path.join(ruta_material, f)) and not f.startswith(".")
    )

    # Compatibilidad hacia atrás: si todavía existen subcarpetas sueltas
    # (por ejemplo restos de la estructura anterior unidad_XX), también
    # se leen sus archivos — ya no se usan para decidir la unidad, pero
    # el contenido no se pierde: entra al mismo pool y el clasificador
    # lo ubica solo, por contenido.
    for entrada in sorted(os.listdir(ruta_material)):
        ruta_entrada = os.path.join(ruta_material, entrada)
        if os.path.isdir(ruta_entrada) and not entrada.startswith("."):
            for f in sorted(os.listdir(ruta_entrada)):
                ruta_completa = os.path.join(ruta_entrada, f)
                if os.path.isfile(ruta_completa):
                    archivos.append(os.path.join(entrada, f))

    if not archivos:
        resultado["advertencias_globales"].append("La carpeta de material está vacía.")
        print("  [!] No se encontraron archivos en la carpeta de material.")
        return resultado

    print(f"Procesando {len(archivos)} archivo(s) en '{ruta_material}' (sin estructura de carpetas por unidad)...")
    for archivo in archivos:
        ruta_completa = os.path.join(ruta_material, archivo)
        documento = extraer_documento(ruta_completa)
        resultado["documentos"].append(documento)

    utilizables = sum(1 for d in resultado["documentos"] if d["utilizable"])
    print(f"\nResumen: {utilizables}/{len(resultado['documentos'])} archivo(s) con contenido utilizable.")

    return resultado


# ==========================================
# BLOQUE DE PRUEBA LOCAL
# ==========================================
if __name__ == "__main__":
    ruta_material = os.path.join(os.path.dirname(__file__), "..", "data", "material")
    material = procesar_material_curso(ruta_material)

    print("\n--- RESUMEN DE EXTRACCIÓN ---")
    for d in material["documentos"]:
        estado = "OK" if d["utilizable"] else "OMITIDO"
        print(f"[{estado}] {d['nombre']} ({d['tipo']}) — {d['num_caracteres']} caracteres")
        for adv in d["advertencias"]:
            print(f"    ! {adv}")
